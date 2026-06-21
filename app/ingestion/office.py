from pathlib import Path

from app.ingestion.errors import ExtractionError
from app.ingestion.types import ExtractedTextBlock, InputType, NormalizedDocument


def extract_docx_file(path: Path, title: str) -> NormalizedDocument:
    blocks = extract_docx_blocks(path)

    if not blocks:
        raise ExtractionError("DOCX does not contain extractable text.")

    return NormalizedDocument(
        title=title,
        source_type=InputType.DOCX,
        blocks=blocks,
    )


def extract_pptx_file(path: Path, title: str) -> NormalizedDocument:
    blocks = extract_pptx_blocks(path)

    if not blocks:
        raise ExtractionError("PPTX does not contain extractable text.")

    return NormalizedDocument(
        title=title,
        source_type=InputType.PPTX,
        blocks=blocks,
    )


def extract_docx_blocks(file_path: Path) -> list[ExtractedTextBlock]:
    try:
        from docx import Document
    except ImportError as exc:
        raise ExtractionError("DOCX ingestion requires python-docx.") from exc

    try:
        document = Document(str(file_path))
    except Exception as exc:
        raise ExtractionError(f"Could not open DOCX file: {exc}") from exc

    blocks: list[ExtractedTextBlock] = []
    current_section: str | None = None

    for paragraph_index, paragraph in enumerate(document.paragraphs):
        text = _clean_text(paragraph.text)

        if not text:
            continue

        style_name = paragraph.style.name if paragraph.style is not None else None
        is_heading = bool(style_name and style_name.lower().startswith("heading"))

        if is_heading:
            current_section = text

        blocks.append(
            ExtractedTextBlock(
                text=text,
                source_page=None,
                source_start_offset=0,
                source_end_offset=len(text),
                metadata={
                    "filename": file_path.name,
                    "source_type": InputType.DOCX.value,
                    "block_type": "paragraph",
                    "paragraph_index": paragraph_index,
                    "style_name": style_name,
                    "section": current_section,
                    "is_heading": is_heading,
                },
            )
        )

    for table_index, table in enumerate(document.tables, start=1):
        table_block = _docx_table_block(
            table=table,
            filename=file_path.name,
            table_index=table_index,
            section=current_section,
        )

        if table_block is not None:
            blocks.append(table_block)

    return blocks


def extract_pptx_blocks(file_path: Path) -> list[ExtractedTextBlock]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractionError("PPTX ingestion requires python-pptx.") from exc

    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise ExtractionError(f"Could not open PPTX file: {exc}") from exc

    blocks: list[ExtractedTextBlock] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_title = _slide_title(slide)
        content_lines: list[str] = []
        table_index = 0

        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                table_index += 1
                table_block = _pptx_table_block(
                    table=shape.table,
                    filename=file_path.name,
                    slide_number=slide_index,
                    slide_title=slide_title,
                    table_index=table_index,
                )

                if table_block is not None:
                    blocks.append(table_block)
                continue

            if not getattr(shape, "has_text_frame", False):
                continue

            text = _clean_text(shape.text)

            if not text or text == slide_title:
                continue

            content_lines.extend(_text_lines(text))

        if slide_title or content_lines:
            text = _slide_text(
                slide_number=slide_index,
                slide_title=slide_title,
                content_lines=content_lines,
            )
            blocks.append(
                ExtractedTextBlock(
                    text=text,
                    source_page=None,
                    source_start_offset=0,
                    source_end_offset=len(text),
                    metadata={
                        "filename": file_path.name,
                        "source_type": InputType.PPTX.value,
                        "block_type": "slide",
                        "slide_number": slide_index,
                        "slide_title": slide_title,
                    },
                )
            )

    return blocks


def _docx_table_block(
    table,
    filename: str,
    table_index: int,
    section: str | None,
) -> ExtractedTextBlock | None:
    rows = _table_rows_from_docx(table)
    text = _table_text(
        title=f"Table {table_index}",
        rows=rows,
    )

    if text is None:
        return None

    headers, data_rows = _split_table_header(rows)
    row_start = 2 if data_rows else None
    row_end = (len(data_rows) + 1) if data_rows else None

    return ExtractedTextBlock(
        text=text,
        source_page=None,
        source_start_offset=0,
        source_end_offset=len(text),
        metadata={
            "filename": filename,
            "source_type": InputType.DOCX.value,
            "block_type": "table",
            "table_index": table_index,
            "section": section,
            "row_start": row_start,
            "row_end": row_end,
            "column_names": headers,
        },
    )


def _pptx_table_block(
    table,
    filename: str,
    slide_number: int,
    slide_title: str | None,
    table_index: int,
) -> ExtractedTextBlock | None:
    rows = _table_rows_from_pptx(table)
    text = _table_text(
        title=f"Slide {slide_number}: {slide_title or 'Untitled'}\n\nTable {table_index}",
        rows=rows,
    )

    if text is None:
        return None

    headers, data_rows = _split_table_header(rows)
    row_start = 2 if data_rows else None
    row_end = (len(data_rows) + 1) if data_rows else None

    return ExtractedTextBlock(
        text=text,
        source_page=None,
        source_start_offset=0,
        source_end_offset=len(text),
        metadata={
            "filename": filename,
            "source_type": InputType.PPTX.value,
            "block_type": "table",
            "slide_number": slide_number,
            "slide_title": slide_title,
            "table_index": table_index,
            "row_start": row_start,
            "row_end": row_end,
            "column_names": headers,
        },
    )


def _slide_title(slide) -> str | None:
    title_shape = getattr(slide.shapes, "title", None)

    if title_shape is None or not getattr(title_shape, "has_text_frame", False):
        return None

    title = _clean_text(title_shape.text)
    return title or None


def _slide_text(
    slide_number: int,
    slide_title: str | None,
    content_lines: list[str],
) -> str:
    heading = f"Slide {slide_number}: {slide_title or 'Untitled'}"
    lines = [heading]

    if slide_title:
        lines.extend(["", f"Title: {slide_title}"])

    if content_lines:
        lines.extend(["Content:", *[f"- {line}" for line in content_lines]])

    return "\n".join(lines)


def _table_text(
    title: str,
    rows: list[list[str]],
) -> str | None:
    headers, data_rows = _split_table_header(rows)

    if not headers and not data_rows:
        return None

    lines = [title]

    if headers:
        lines.append(f"Headers: {' | '.join(headers)}")

    for row_index, row in enumerate(data_rows, start=1):
        lines.append(f"Row {row_index}: {' | '.join(row)}")

    return "\n".join(lines)


def _split_table_header(rows: list[list[str]]) -> tuple[list[str], list[list[str]]]:
    if not rows:
        return [], []

    headers = _normalize_headers(rows[0])
    data_rows = [row for row in rows[1:] if any(cell.strip() for cell in row)]

    return headers, data_rows


def _normalize_headers(values: list[str]) -> list[str]:
    headers: list[str] = []

    for index, value in enumerate(values, start=1):
        header = _clean_text(value)
        headers.append(header or f"Column {index}")

    return headers


def _table_rows_from_docx(table) -> list[list[str]]:
    rows: list[list[str]] = []

    for row in table.rows:
        values = [_clean_text(cell.text) for cell in row.cells]
        if any(values):
            rows.append(values)

    return rows


def _table_rows_from_pptx(table) -> list[list[str]]:
    rows: list[list[str]] = []

    for row in table.rows:
        values = [_clean_text(cell.text) for cell in row.cells]
        if any(values):
            rows.append(values)

    return rows


def _clean_text(value: str | None) -> str:
    if value is None:
        return ""

    return " ".join(str(value).split()).strip()


def _text_lines(value: str) -> list[str]:
    return [line for line in (_clean_text(line) for line in value.splitlines()) if line]
