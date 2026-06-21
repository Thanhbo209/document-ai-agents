from pathlib import Path

from fastapi.testclient import TestClient
from sqlalchemy.orm import Session

from app.ingestion.loader import load_document
from app.ingestion.office import extract_docx_blocks, extract_pptx_blocks
from app.ingestion.types import InputType
from app.repositories.documents import DocumentRepository
from tests.helpers import create_authenticated_workspace


def create_sample_docx(path: Path) -> None:
    from docx import Document

    document = Document()
    document.add_heading("Employee Directory", level=1)
    document.add_paragraph("The staffing plan includes AI and engineering roles.")

    table = document.add_table(rows=1, cols=3)
    headers = table.rows[0].cells
    headers[0].text = "Name"
    headers[1].text = "Role"
    headers[2].text = "Start Date"

    row = table.add_row().cells
    row[0].text = "Thanh"
    row[1].text = "AI Intern"
    row[2].text = "2026-06-01"

    row = table.add_row().cells
    row[0].text = "Alex"
    row[1].text = "Engineer"
    row[2].text = "2026-06-02"

    document.save(path)


def create_sample_pptx(path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    title_layout = presentation.slide_layouts[5]

    slide = presentation.slides.add_slide(title_layout)
    slide.shapes.title.text = "Product Roadmap"

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5), Inches(1.5))
    text_box.text_frame.text = "Phase 1: Upload\nPhase 2: Retrieval\nPhase 3: Review"

    table_shape = slide.shapes.add_table(
        3,
        2,
        Inches(1),
        Inches(3.2),
        Inches(5),
        Inches(1.3),
    )
    table = table_shape.table
    table.cell(0, 0).text = "Phase"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Upload"
    table.cell(1, 1).text = "Done"
    table.cell(2, 0).text = "Retrieval"
    table.cell(2, 1).text = "Planned"

    slide = presentation.slides.add_slide(title_layout)
    slide.shapes.title.text = "Launch Notes"
    slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(5), Inches(1)
    ).text = "Launch depends on review workflow readiness."

    presentation.save(path)


def test_extract_docx_blocks_extracts_paragraphs_headings_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "sample.docx"
    create_sample_docx(path)

    blocks = extract_docx_blocks(path)
    paragraph_blocks = [block for block in blocks if block.metadata["block_type"] == "paragraph"]
    table_blocks = [block for block in blocks if block.metadata["block_type"] == "table"]

    assert any(block.text == "Employee Directory" for block in paragraph_blocks)
    assert any("staffing plan" in block.text for block in paragraph_blocks)
    assert len(table_blocks) == 1

    table = table_blocks[0]
    assert "Headers: Name | Role | Start Date" in table.text
    assert "Row 1: Thanh | AI Intern | 2026-06-01" in table.text
    assert table.metadata["source_type"] == "docx"
    assert table.metadata["table_index"] == 1
    assert table.metadata["row_start"] == 2
    assert table.metadata["row_end"] == 3
    assert table.metadata["column_names"] == ["Name", "Role", "Start Date"]


def test_extract_pptx_blocks_preserves_slide_titles_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "sample.pptx"
    create_sample_pptx(path)

    blocks = extract_pptx_blocks(path)
    slide_blocks = [block for block in blocks if block.metadata["block_type"] == "slide"]
    table_blocks = [block for block in blocks if block.metadata["block_type"] == "table"]

    assert len(slide_blocks) == 2
    assert "Slide 1: Product Roadmap" in slide_blocks[0].text
    assert "- Phase 1: Upload" in slide_blocks[0].text
    assert slide_blocks[0].metadata["slide_number"] == 1
    assert slide_blocks[0].metadata["slide_title"] == "Product Roadmap"

    assert len(table_blocks) == 1
    table = table_blocks[0]
    assert "Headers: Phase | Status" in table.text
    assert table.metadata["source_type"] == "pptx"
    assert table.metadata["slide_number"] == 1
    assert table.metadata["slide_title"] == "Product Roadmap"
    assert table.metadata["column_names"] == ["Phase", "Status"]


def test_load_document_routes_docx_and_pptx(tmp_path: Path) -> None:
    docx_path = tmp_path / "sample.docx"
    pptx_path = tmp_path / "sample.pptx"
    create_sample_docx(docx_path)
    create_sample_pptx(pptx_path)

    docx_document = load_document(docx_path, InputType.DOCX, "sample")
    pptx_document = load_document(pptx_path, InputType.PPTX, "sample")

    assert docx_document.source_type == InputType.DOCX
    assert any(block.metadata["block_type"] == "table" for block in docx_document.blocks)
    assert pptx_document.source_type == InputType.PPTX
    assert any(block.metadata["slide_number"] == 1 for block in pptx_document.blocks)


def test_upload_docx_and_pptx_succeeds(
    client: TestClient,
    db_session: Session,
    tmp_path: Path,
) -> None:
    workspace_id, headers = create_authenticated_workspace(
        db_session,
        email="office-upload-owner@example.com",
    )
    document_repo = DocumentRepository(db_session)

    docx_path = tmp_path / "upload.docx"
    pptx_path = tmp_path / "upload.pptx"
    create_sample_docx(docx_path)
    create_sample_pptx(pptx_path)

    for filename, content_type, path in [
        (
            "upload.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            docx_path,
        ),
        (
            "upload.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            pptx_path,
        ),
    ]:
        response = client.post(
            f"/api/v1/workspaces/{workspace_id}/documents/upload",
            headers=headers,
            files={
                "file": (
                    filename,
                    path.read_bytes(),
                    content_type,
                )
            },
        )

        assert response.status_code == 201
        payload = response.json()
        assert payload["status"] == "succeeded"
        assert payload["chunks_created"] >= 1

        chunks = document_repo.list_chunks_for_document(
            workspace_id=workspace_id,
            document_id=payload["document_id"],
        )
        assert chunks
        assert chunks[0].source_metadata["source_type"] in {"docx", "pptx"}
